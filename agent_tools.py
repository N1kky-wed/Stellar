import time
import uuid
import json
import base64
import os
import re

# Heavy imports (tavily, genai, types, asyncio) have been removed from the global scope and are loaded lazily.

import logging
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

from typing import List, Optional
from flask import g, has_request_context

def request_user_interaction(html_ui: str, goal: str, status: str, timeout: int) -> str:
    """
    Renders an interactive UI on the user's screen and pauses your execution until the user interacts with it to achieve the goal. 
    Use this when you need the user to play a minigame, fill out a custom form, or provide structured visual input.
    The `html_ui` must contain Javascript that eventually calls `window.stellar.finish(data)` with a JSON object or string containing the result.
    The result will be returned to you directly.
    """
    pass # This is intercepted by gemini_generate in app.py

def _get_effective_session():
    """Helper to get session data safely in both request and background thread contexts."""
    from flask import session, has_app_context
    if has_request_context():
        return session
    
    class SafeSession(dict):
        """
        A dict-like session wrapper that acts as a safe fallback when operating outside
        a traditional Flask request context (e.g. in background thread execution).
        """
        def __init__(self):
            """
            Initialize SafeSession by pre-populating session keys from the global thread context.
            """
            super().__init__()
            # Pre-populate from g which is set in background_thread_runner
            # Only access g if we have an app context
            if has_app_context():
                self['user_id'] = getattr(g, 'user_id', None)
                self['username'] = getattr(g, 'username', None)
                self['current_chat_id'] = getattr(g, 'chat_id', None)
                # Compatibility for session.sid and session.modified
                self.sid = getattr(g, 'session_id', 'no_session')
            else:
                self.sid = 'no_session'
            self.modified = False
    return SafeSession()

def web_search(
    action: str,
    status: str,
    timeout: int,
    query: Optional[str] = None,
    url: Optional[str] = None,
    urls: Optional[List[str]] = None,
    topic: str = "general",  # 'general', 'news', 'finance'
    search_depth: str = "advanced",  # 'ultra-fast', 'fast', 'basic', 'advanced'
    extract_depth: str = "basic",  # 'basic', 'advanced'
    auto_parameters: bool = False,
    exact_match: bool = False,
    time_range: Optional[str] = None,  # 'd', 'w', 'm', 'y'
    start_date: Optional[str] = None,  # YYYY-MM-DD
    end_date: Optional[str] = None,  # YYYY-MM-DD
    max_results: int = 5,
    chunks_per_source: int = 3,
    include_answer: bool = True,
    include_raw_content: bool = False,
    include_images: bool = True,
    include_image_descriptions: bool = True,
    include_favicon: bool = False,
    include_usage: bool = False,
    include_domains: Optional[List[str]] = None,
    exclude_domains: Optional[List[str]] = None,
    select_domains: Optional[List[str]] = None,
    select_paths: Optional[List[str]] = None,
    exclude_paths: Optional[List[str]] = None,
    country: Optional[str] = None,
    format: str = "markdown",  # 'markdown', 'text'
    instructions: Optional[str] = None,
    max_depth: int = 2,
    max_breadth: int = 20,
    limit: int = 50,
    allow_external: bool = True
) -> str:
    """Unified Web Search, Extraction, Crawling, and Mapping Tool.
    
    Args:
        action: 'tavily_search', 'tavily_extract', 'tavily_crawl', 'tavily_map'.
        status: Status update for the user.
        query: Search term or semantic intent.
        url: Root URL for crawl/map.
        urls: List of URLs for extraction (max 20).
        topic: Category of search ('general', 'news', 'finance').
        search_depth: 'ultra-fast' (instant), 'fast' (quick snippets), 'basic' (high relevance), 'advanced' (deep analysis).
        extract_depth: 'basic' or 'advanced' (extracts tables/embedded media).
        auto_parameters: Let Tavily optimize parameters based on intent.
        exact_match: Ensure exact quoted phrases bypass synonyms.
        time_range: Relative time ('day', 'week', 'month', 'year' or shorthand 'd', 'w', 'm', 'y').
        start_date / end_date: Specific dates formatted as YYYY-MM-DD.
        max_results: Number of search results (0-20).
        chunks_per_source: Max snippets per source (1-5).
        include_answer: AI generated summary.
        include_raw_content: Returns parsed HTML text/markdown.
        include_images: Extract images from search/URLs/Crawl.
        include_image_descriptions: Add LLM descriptions to images.
        include_favicon: Include favicon URL.
        include_usage: Return API credit usage stats.
        include_domains / exclude_domains: For Search/Crawl limiting.
        select_domains / select_paths / exclude_paths: Regex patterns for Map/Crawl limits.
        country: Boost search results from specific country.
        format: Extraction output ('markdown' or 'text').
        instructions: Natural language guidance for crawlers.
        max_depth: Link click depth for Map/Crawl.
        max_breadth: Max links to follow per level for Map/Crawl.
        limit: Max pages to process for Map/Crawl.
        allow_external: Follow links to external domains during Map/Crawl.
        timeout: Wait time in seconds before failing.
    """
    try:
        from app import TAVILY_API_KEY, TAVILY_BACKUP_API_KEYS
        tavily_keys = [TAVILY_API_KEY] + [bk for bk in TAVILY_BACKUP_API_KEYS if bk]
        tavily_keys = [k for k in dict.fromkeys(tavily_keys) if k]

        def execute_tavily_with_retries(operation, **kwargs):
            """
            Execute a Tavily API operation with automatic key rotation and retry logic.

            Args:
                operation (str): The Tavily operation to execute ('search', 'extract', 'crawl', or 'map').
                **kwargs: Keyword arguments passed to the Tavily SDK.

            Returns:
                dict: The response returned by the Tavily API.
            """
            if not tavily_keys:
                raise Exception("Tavily search failed: API Key missing.")
            last_error = None
            for key in tavily_keys:
                try:
                    t_client = TavilyClient(api_key=key)
                    t0 = time.time()
                    if operation == 'search':
                        res = t_client.search(**kwargs)
                    elif operation == 'extract':
                        res = t_client.extract(**kwargs)
                    elif operation == 'crawl':
                        res = t_client.crawl(**kwargs)
                    elif operation == 'map':
                        res = t_client.map(**kwargs)
                    logger.info("Tavily %s call completed duration_sec=%.2f", operation, time.time() - t0)
                    return res
                except Exception as e:
                    logger.error("Error in Tavily operation=%s key_suffix=%s error=%s", operation, key[-4:] if key else '', e, exc_info=True)
                    last_error = e
                    continue
            raise Exception(f"All Tavily API keys exhausted. Last error: {str(last_error)}")

        # Strip None values out of kwargs to prevent SDK validation errors
        def clean_kwargs(kwargs_dict):
            """
            Strip out all None values from a dictionary.

            Args:
                kwargs_dict (dict): The dictionary to clean.

            Returns:
                dict: The cleaned dictionary.
            """
            return {k: v for k, v in kwargs_dict.items() if v is not None}

        # Deterministic Image Verification Filter
        def _verify_tavily_images(data):
            """
            Verify that images returned by Tavily are actually online.

            Args:
                data (dict): The raw Tavily search/extract result.

            Returns:
                dict: The results dictionary containing verified images.
            """
            if not include_images: return data
            import requests
            from concurrent.futures import ThreadPoolExecutor
            def _is_online(url):
                """
                Check if the given image URL returns a 200 HTTP status code.

                Args:
                    url (str): The image URL.

                Returns:
                    bool: True if online; False otherwise.
                """
                try:
                    r = requests.get(url, stream=True, timeout=3)
                    r.close()
                    return r.status_code == 200
                except Exception as err:
                    logger.debug("Tavily image URL %s offline: %s", url, err)
                    return False
            
            urls = set()
            for i in data.get('images', []):
                urls.add(i if isinstance(i, str) else i.get('url'))
            for r in data.get('results', []):
                for i in r.get('images', []):
                    urls.add(i if isinstance(i, str) else i.get('url'))
            urls.discard(None)
            
            if not urls: return data
            
            with ThreadPoolExecutor(max_workers=10) as ex:
                valid_urls = {u for u, ok in zip(urls, ex.map(_is_online, urls)) if ok}
                
            if 'images' in data:
                data['images'] = [i for i in data['images'] if (i if isinstance(i, str) else i.get('url')) in valid_urls]
            for r in data.get('results', []):
                if 'images' in r:
                    r['images'] = [i for i in r['images'] if (i if isinstance(i, str) else i.get('url')) in valid_urls]
            return data

        # 2. Tavily Search
        if action == "tavily_search":
            if not query: return json.dumps({"error": "Missing 'query' for tavily_search"})
            kwargs = clean_kwargs({
                "query": query, "topic": topic, "search_depth": search_depth,
                "auto_parameters": auto_parameters, "max_results": max_results,
                "time_range": time_range, "start_date": start_date, "end_date": end_date,
                "include_answer": include_answer, "include_raw_content": include_raw_content,
                "include_images": include_images, "include_image_descriptions": include_image_descriptions,
                "include_domains": include_domains, "exclude_domains": exclude_domains,
                "country": country, "exact_match": exact_match, 
                "include_favicon": include_favicon, "include_usage": include_usage,
                "timeout": timeout
            })
            if search_depth in ["advanced", "fast"]:
                kwargs["chunks_per_source"] = chunks_per_source
            
            res = execute_tavily_with_retries('search', **kwargs)
            return json.dumps({"tool": "tavily_search", "data": _verify_tavily_images(res)}, indent=2)

        # 3. Tavily Extract
        elif action == "tavily_extract":
            if not urls: return json.dumps({"error": "Missing 'urls' list for tavily_extract"})
            kwargs = clean_kwargs({
                "urls": urls, "extract_depth": extract_depth, "format": format,
                "include_images": include_images, "include_favicon": include_favicon,
                "include_usage": include_usage, "timeout": timeout
            })
            if query:
                kwargs["query"] = query
                kwargs["chunks_per_source"] = chunks_per_source
                
            res = execute_tavily_with_retries('extract', **kwargs)
            return json.dumps({"tool": "tavily_extract", "data": _verify_tavily_images(res)}, indent=2)

        # 4. Tavily Crawl
        elif action == "tavily_crawl":
            if not url: return json.dumps({"error": "Missing 'url' for tavily_crawl"})
            kwargs = clean_kwargs({
                "url": url, "max_depth": max_depth, "max_breadth": max_breadth, 
                "limit": limit, "instructions": instructions, "select_paths": select_paths,
                "select_domains": select_domains, "exclude_paths": exclude_paths, 
                "exclude_domains": exclude_domains, "allow_external": allow_external,
                "include_images": include_images, "extract_depth": extract_depth,
                "format": format, "include_favicon": include_favicon, 
                "include_usage": include_usage, "timeout": timeout,
                "chunks_per_source": chunks_per_source
            })
            res = execute_tavily_with_retries('crawl', **kwargs)
            return json.dumps({"tool": "tavily_crawl", "data": _verify_tavily_images(res)}, indent=2)

        # 5. Tavily Map
        elif action == "tavily_map":
            if not url: return json.dumps({"error": "Missing 'url' for tavily_map"})
            kwargs = clean_kwargs({
                "url": url, "max_depth": max_depth, "max_breadth": max_breadth, 
                "limit": limit, "instructions": instructions, "select_paths": select_paths,
                "select_domains": select_domains, "exclude_paths": exclude_paths, 
                "exclude_domains": exclude_domains, "allow_external": allow_external,
                "include_usage": include_usage, "timeout": timeout
            })
            # Map doesn't return images typically, but we wrap it just in case it ever does
            res = execute_tavily_with_retries('map', **kwargs)
            return json.dumps({"tool": "tavily_map", "data": _verify_tavily_images(res)}, indent=2)

        else:
            return json.dumps({"error": f"Unknown action: '{action}'."})

    except Exception as e:
        logger.exception("Error caught: %s", e)
        return json.dumps({"error": str(e)})

def send_self_email(subject: str, body: str, status: str, timeout: int, attachment_path: str = None) -> str:
    """Secure Closed-Loop Mailer. Sends reports/files ONLY to your registered email address.
    Args:
        subject: Subject line.
        body: Content of the email.
        status: Status update for the user.
        attachment_path: Optional path to a file in /outputs or /uploads to attach.
    """
    session = _get_effective_session()
    import smtplib, mimetypes, os, sqlite3, markdown
    from email.message import EmailMessage

    # Handle background context where session is unavailable
    user_email = None
    from flask import has_app_context
    if has_request_context():
        user_email = session.get('username')
        
    if not user_email and has_app_context() and hasattr(g, 'user_id'):
        try:
            from app import DATABASE_NAME
            conn = sqlite3.connect(DATABASE_NAME)
            conn.execute("PRAGMA journal_mode=WAL;")
            conn.execute("PRAGMA busy_timeout=5000;")
            conn.row_factory = sqlite3.Row
            user = conn.execute('SELECT username FROM users WHERE id = ?', (g.user_id,)).fetchone()
            if user:
                user_email = user['username']
            conn.close()
        except Exception as db_e:
            logger.error("Error fetching email for background task error=%s", db_e, exc_info=True)

    if not user_email:
        return "Error: Could not determine recipient email address (Session/Context missing)."

    sender_email = os.getenv("EMAIL_USER")
    sender_password = os.getenv("EMAIL_PASS")

    msg = EmailMessage()
    msg['Subject'] = f"[STELLAR] {subject}"
    msg['From'] = f"Stellar System <{sender_email}>"
    msg['To'] = user_email
    
    # Text Version
    text_content = f"{body}\n\n---\nTransmission from your Stellar Environment."
    msg.set_content(text_content)

    # HTML Version (Markdown Rendered)
    try:
        html_body = markdown.markdown(body, extensions=['extra', 'codehilite', 'tables'])
        html_content = f"""
        <html>
          <head>
            <style>
              body {{ font-family: sans-serif; line-height: 1.6; color: #333; }}
              code {{ background: #f4f4f4; padding: 2px 4px; border-radius: 4px; }}
              pre {{ background: #f4f4f4; padding: 10px; border-radius: 4px; overflow-x: auto; }}
              table {{ border-collapse: collapse; width: 100%; margin: 10px 0; }}
              th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
              th {{ background-color: #f2f2f2; }}
              hr {{ border: 0; border-top: 1px solid #eee; margin: 20px 0; }}
              .footer {{ font-size: 0.85em; color: #777; margin-top: 20px; }}
            </style>
          </head>
          <body>
            {html_body}
            <div class="footer">
              <hr>
              Transmission from your Stellar Environment.
            </div>
          </body>
        </html>
        """
        msg.add_alternative(html_content, subtype='html')
    except Exception as md_e:
        logger.error("Markdown rendering failed error=%s", md_e, exc_info=True)

    # Robust Attachment Handling
    if attachment_path:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        allowed_dirs = [
            os.path.abspath(os.path.join(base_dir, 'outputs')),
            os.path.abspath(os.path.join(base_dir, 'uploads')),
            os.path.abspath(os.path.join(base_dir, 'sandbox_runs'))
        ]
        
        candidate_path = os.path.abspath(attachment_path) if os.path.isabs(attachment_path) else os.path.abspath(os.path.join(base_dir, attachment_path))
        is_allowed = any(candidate_path.startswith(d + os.sep) for d in allowed_dirs)
        
        # Fallback: check if the exact filename exists in outputs or uploads
        if not is_allowed or not os.path.exists(candidate_path):
            filename = os.path.basename(attachment_path)
            for d in [allowed_dirs[0], allowed_dirs[1]]: # outputs and uploads
                fallback_path = os.path.join(d, filename)
                if os.path.exists(fallback_path) and os.path.isfile(fallback_path):
                    candidate_path = fallback_path
                    is_allowed = True
                    break

        resolved_path = candidate_path if (is_allowed and os.path.exists(candidate_path) and os.path.isfile(candidate_path)) else None

        if resolved_path:
            ctype, encoding = mimetypes.guess_type(resolved_path)
            if ctype is None or encoding is not None:
                ctype = 'application/octet-stream'
            maintype, subtype = ctype.split('/', 1)
            
            try:
                with open(resolved_path, 'rb') as fp:
                    msg.add_attachment(
                        fp.read(),
                        maintype=maintype,
                        subtype=subtype,
                        filename=os.path.basename(resolved_path)
                    )
                logger.info(f"Successfully attached file: {resolved_path}")
            except Exception as att_e:
                logger.error("Failed to attach file path=%s error=%s", resolved_path, att_e, exc_info=True)
        else:
            logger.warning("Attachment path invalid denied or not found path=%s", attachment_path)

    try:
        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
            smtp.login(sender_email, sender_password)
            smtp.send_message(msg)
        return f"Success: Email sent to {user_email}."
    except Exception as e:
        logger.error("Mail Failure error=%s", e, exc_info=True)
        return f"Mail Failure: {str(e)}"

def schedule_task(task_prompt: str, status: str, timeout: int, action: str = "schedule", task_id: int = None, execute_at: str = None, recurring_minutes: int = 0, metadata: str = None) -> str:
    """Schedules, lists, or cancels autonomous tasks.
    Args:
        task_prompt: Instructions for the AI to follow (required for 'schedule').
        status: Status update for the user.
        action: 'schedule' (default), 'list' (to see pending tasks), 'cancel' (to stop a task), or 'edit' (to modify a task).
        task_id: The ID of the task to cancel or edit (required for 'cancel'/'edit').
        execute_at: ISO datetime 'YYYY-MM-DD HH:MM:SS' for one-time tasks.
        recurring_minutes: Minutes between executions for repeating tasks.
        metadata: Optional scratchpad for task-specific state (retry counts, transient notes).
    """
    from flask import request
    session = _get_effective_session()
    from app import get_db
    import sqlite3

    db = get_db()
    u_id = None
    c_id = None
    # Server-side model detection to prevent hallucination
    current_model = getattr(g, 'model_id', 'gemma-4-31b-it')
    
    if has_request_context():
        u_id = session.get('user_id')
        c_id = session.get('current_chat_id')
        if request.is_json:
            current_model = request.json.get('model_id', current_model)
        
    if not u_id: u_id = getattr(g, 'user_id', None)
    if not c_id: c_id = getattr(g, 'chat_id', None)
    
    if not u_id or not c_id:
        return "Error: Could not determine User or Chat ID for scheduling."

    if action == "list":
        try:
            cursor = db.execute('SELECT id, task_prompt, model_id, execute_at, recurring_minutes, metadata FROM scheduled_tasks WHERE user_id = ? AND is_active = 1', (u_id,))
            tasks = cursor.fetchall()
            if not tasks: return "No active scheduled tasks found."
            output = "### Active Scheduled Tasks:\n"
            for t in tasks:
                meta_snip = f" | State: {t['metadata'][:30]}..." if t['metadata'] else ""
                output += f"- **ID {t['id']}** [{t['model_id']}]: \"{t['task_prompt'][:50]}...\"{meta_snip} | Next: {t['execute_at']} | Every: {t['recurring_minutes']}m\n"
            return output
        except Exception as e:
            logger.exception("Error caught: %s", e)
            return f"Error listing tasks: {str(e)}"

    elif action == "cancel":
        if not task_id: return "Error: 'task_id' is required to cancel a task."
        cursor = db.execute('UPDATE scheduled_tasks SET is_active = 0 WHERE id = ? AND user_id = ?', (task_id, u_id))
        db.commit()
        if cursor.rowcount == 0:
            return f"Error: Task {task_id} not found or already inactive."
        return f"Success: Task {task_id} has been cancelled."

    elif action == "edit":
        if not task_id: return "Error: 'task_id' is required to edit a task."
        # ... existing edit logic ...
        updates = []
        params = []
        if task_prompt:
            updates.append("task_prompt = ?")
            params.append(task_prompt)
        if execute_at:
            updates.append("execute_at = ?")
            params.append(execute_at)
        if recurring_minutes is not None:
            updates.append("recurring_minutes = ?")
            params.append(recurring_minutes)
        if metadata:
            updates.append("metadata = ?")
            params.append(metadata)
        
        if not updates:
            return "Error: No parameters provided to edit."
        
        params.append(task_id)
        params.append(u_id)
        cursor = db.execute(f'UPDATE scheduled_tasks SET {", ".join(updates)} WHERE id = ? AND user_id = ?', tuple(params))
        db.commit()
        if cursor.rowcount == 0:
            return f"Error: Task {task_id} not found or already inactive."
        return f"Success: Task {task_id} has been updated."

    # Default: Schedule
    cursor = db.execute('SELECT COUNT(*) FROM scheduled_tasks WHERE user_id = ? AND is_active = 1', (u_id,))
    if cursor.fetchone()[0] >= 10:
        return "Error: Maximum number of active scheduled tasks (10) reached. Please cancel some tasks before scheduling more."

    cursor = db.execute('INSERT INTO scheduled_tasks (user_id, chat_id, task_prompt, model_id, execute_at, recurring_minutes, metadata) VALUES (?,?,?,?,?,?,?)',
               (u_id, c_id, task_prompt, current_model, execute_at, recurring_minutes, metadata))
    new_id = cursor.lastrowid
    db.commit()
    logger.info("Task scheduled task_id=%s user_id=%s chat_id=%s model_id=%s recurring_minutes=%s", new_id, u_id, c_id, current_model, recurring_minutes)
    return f"Task scheduled (ID: {new_id})! {current_model} is locked for this persistent automation."

def generate_image(model: str, prompt: str, status: str, timeout: int, quality: str = "1K", aspect_ratio: str = "1:1", reference_images: list[str] = None) -> str:
    """Generates an image using Gemini's Imagen model.
    Args:
        model: 'gemini-3.1-flash-image-preview' or 'gemini-3-pro-image-preview'
        prompt: detailed descriptive prompt for the image
        status: Status update for the user.
        quality: Supported tiers are "512", "1K", "2K", "4K". (Default: "1K")
        aspect_ratio: Supported ratios: '1:1', '3:4', '4:3', '9:16', '16:9'.
        reference_images: List of filenames from the chat context to use as reference/conditioning (up to 14).
    """
    from app import PRIMARY_API_KEY, UPLOAD_FOLDER, BACKUP_API_KEYS, KEY_MANAGER, parse_quota_block_duration
    session = _get_effective_session()
    import os
    import mimetypes
    import uuid
    from google import genai
    from google.genai import types
    
    raw_keys = [PRIMARY_API_KEY] + [bk for bk in BACKUP_API_KEYS if bk]
    keys_to_try = [k for k in dict.fromkeys(raw_keys) if k]
    
    active_keys = [k for k in keys_to_try if not KEY_MANAGER.is_key_blocked(k, model)[0]]
    if not active_keys:
        active_keys = keys_to_try
    keys_to_try = active_keys
    
    # Ensure aspect_ratio is one of the strictly supported API values
    valid_ratios = ["1:1", "3:4", "4:3", "9:16", "16:9"]
    if aspect_ratio not in valid_ratios:
        aspect_ratio = "1:1"
        
    quality_lower = quality.lower()
    if "4k" in quality_lower or "high" in quality_lower or "hd" in quality_lower:
        img_size = "4K"
    elif "2k" in quality_lower:
        img_size = "2K"
    elif "1k" in quality_lower or "standard" in quality_lower:
        img_size = "1K"
    elif "512" in quality_lower:
        img_size = "512"
    else:
        img_size = "1K"
    
    image_config_args = {"aspect_ratio": aspect_ratio}
    try:
        if hasattr(types, 'ImageConfig'):
            valid_fields = getattr(types.ImageConfig, '__pydantic_fields__', {})
            if 'image_size' in valid_fields:
                image_config_args["image_size"] = img_size
    except Exception:
        pass
    
    # Resolve reference images
    parts = [types.Part.from_text(text=prompt)]
    if reference_images:
        try:
            chat_id = session.get('current_chat_id')
            context_id = str(chat_id) if chat_id else getattr(session, 'sid', 'no_session')
            local_dir = os.path.join(UPLOAD_FOLDER, context_id)
            for img_name in reference_images[:14]:
                img_path = os.path.join(local_dir, img_name)
                if os.path.exists(img_path):
                    mime_type, _ = mimetypes.guess_type(img_path)
                    with open(img_path, "rb") as f:
                        parts.append(types.Part.from_bytes(data=f.read(), mime_type=mime_type or "image/png"))
        except Exception as e:
            logger.error("Error loading reference images error=%s", e, exc_info=True)

    last_error = None
    for current_key in keys_to_try:
        try:
            client = genai.Client(api_key=current_key)
            t0 = time.time()
            response = client.models.generate_content(
                model=model,
                contents=parts,
                config=types.GenerateContentConfig(
                    image_config=image_config_args,
                    response_modalities=["IMAGE"]
                )
            )
            duration = time.time() - t0
            logger.info("Gemini API call completed model=%s duration_sec=%.2f purpose=generate_image", model, duration)
            
            for part in response.candidates[0].content.parts:
                if hasattr(part, 'inline_data') and part.inline_data:
                    img_data = part.inline_data.data
                    
                    output_dir = "outputs"
                    os.makedirs(output_dir, exist_ok=True)
                    
                    mime_type = getattr(part.inline_data, 'mime_type', 'image/png')
                    ext = "png"
                    if "jpeg" in mime_type: ext = "jpg"
                    elif "webp" in mime_type: ext = "webp"
                    
                    filename = f"gen_{uuid.uuid4().hex[:8]}.{ext}"
                    file_path = os.path.join(output_dir, filename)
                    
                    with open(file_path, "wb") as f:
                        f.write(img_data)
                    
                    return f"![Generated Image](https://stellarai.live/view/{filename})"
            return "Error: Image model returned no visual data. Since image generation failed, please use the web search tool to find relevant images on the web instead."
        except Exception as e:
            logger.error("Error in generate_image tool error=%s", e, exc_info=True)
            error_string = str(e).lower()
            if ('429' in error_string or '403' in error_string or '503' in error_string or '500' in error_string or 'resource_exhausted' in error_string or 'quota' in error_string):
                block_duration, block_reason = parse_quota_block_duration(error_string)
                block_scope = None if ('403' in error_string or 'permission_denied' in error_string or 'invalid' in error_string) else model
                KEY_MANAGER.block_key(current_key, block_scope, block_duration, block_reason)
                logger.warning("Globally blocked API key hash=%s block_duration_sec=%d model=%s reason=%s error=generate_image", hash(current_key), block_duration, block_scope, block_reason)
                last_error = e
                continue
            return f"Error generating image: {str(e)}. Since image generation failed, please use the web search tool to find relevant images on the web instead."
    
    return f"Error: All API keys exhausted. Last error: {str(last_error)}. Since image generation failed, please use the web search tool to find relevant images on the web instead."


def report_process_issue(topic: str, issue_description: str, technical_context: str, status: str, timeout: int) -> str:
    """Reports technical bottlenecks, process failures, or feedback on internal tool execution.
    Args:
        topic: A concise label for the issue (e.g., 'SIGKILL', 'Path Alignment', 'Port Latency').
        issue_description: A detailed explanation of what went wrong and how it impacted the task.
        technical_context: Raw logs, error codes, environment details, or reproduction steps.
        status: Status update for the user.
    """
    import sqlite3
    from app import DATABASE_NAME
    from flask import g, has_request_context, session
    
    u_id = None
    c_id = None

    if has_request_context():
        u_id = session.get('user_id')
        c_id = session.get('current_chat_id')

    if not u_id: u_id = getattr(g, 'user_id', None)
    if not c_id: c_id = getattr(g, 'chat_id', None)

    try:
        conn = sqlite3.connect(DATABASE_NAME)
        conn.execute("PRAGMA journal_mode=WAL;")
        conn.execute("PRAGMA busy_timeout=5000;")

        conn.execute('''
            INSERT INTO agent_feedback (user_id, chat_id, topic, issue_description, technical_context, status)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (u_id, c_id, topic, issue_description, technical_context, 'open'))

        conn.commit()
        conn.close()
        logger.info("Agent feedback stored successfully user_id=%s chat_id=%s topic=%s", u_id, c_id, topic)

        import subprocess
        import os
        import sys
        resolver_path = os.path.join(os.path.dirname(__file__), 'issue_resolver.py')
        log_path = os.path.join(os.path.dirname(__file__), 'issue_resolver.log')
        with open(log_path, 'a') as log_file:
            subprocess.Popen([sys.executable, resolver_path], stdout=log_file, stderr=log_file)

        return "Feedback successfully reported and stored for developer review."
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error reporting feedback: {str(e)}"

def compress_memory(target: str, state_document: str, status: str, timeout: int) -> str:
    """Compresses the chat's memory by archiving old tool logs and/or messages.
    Called when the system warns that context window usage is high.
    Preserves a structured state document so critical context is not lost.
    
    Args:
        target: What to compress. Must be one of: 'tool_logs', 'chat_messages', or 'both'.
        state_document: A structured summary preserving current objective, key discoveries, modified files, and pending blockers.
        status: Status update for the user.
        timeout: Execution timeout in seconds.
    """
    import sqlite3
    try:
        from flask import g
        from app import DATABASE_NAME
        
        chat_id = getattr(g, 'chat_id', None)
        if not chat_id:
            return "Error: No active chat session found."
        
        if target not in ['tool_logs', 'chat_messages', 'both']:
            return f"Error: Invalid target '{target}'. Must be 'tool_logs', 'chat_messages', or 'both'."
        
        if not state_document or len(state_document.strip()) < 50:
            return "Error: state_document is too short. You must write a thorough summary of the current state before compressing."
        
        conn = sqlite3.connect(DATABASE_NAME)
        conn.execute("PRAGMA journal_mode=WAL;")
        conn.execute("PRAGMA busy_timeout=5000;")
        conn.row_factory = sqlite3.Row
        tools_archived = 0
        msgs_archived = 0
        
        if target in ['tool_logs', 'both']:
            # Hide all tool calls EXCEPT the 10 most recent
            cursor = conn.execute(
                'SELECT id FROM tool_calls WHERE chat_id = ? AND hidden = 0 ORDER BY id DESC LIMIT 10',
                (chat_id,)
            )
            keep_ids = [row['id'] for row in cursor.fetchall()]
            
            if keep_ids:
                placeholders = ','.join('?' * len(keep_ids))
                cursor = conn.execute(
                    f'UPDATE tool_calls SET hidden = 1 WHERE chat_id = ? AND hidden = 0 AND id NOT IN ({placeholders})',
                    [chat_id] + keep_ids
                )
            else:
                cursor = conn.execute(
                    'UPDATE tool_calls SET hidden = 1 WHERE chat_id = ? AND hidden = 0',
                    (chat_id,)
                )
            tools_archived = cursor.rowcount
        
        if target in ['chat_messages', 'both']:
            cursor = conn.execute(
                'SELECT id, message_type, message_content, timestamp FROM messages WHERE chat_id = ? AND hidden = 0 ORDER BY id ASC',
                (chat_id,)
            )
            all_msgs = cursor.fetchall()
            if len(all_msgs) > 4:
                msgs_to_archive = all_msgs[:-4]
                keep_ids = [row['id'] for row in all_msgs[-4:]]
                
                archived_messages_text = "\n--- COMPRESSION EVENT ---\n"
                for m in msgs_to_archive:
                    role = 'User' if m['message_type'] == 'user' else 'Stellar'
                    archived_messages_text += f"[{m['timestamp']}] {role}:\n{m['message_content']}\n\n"
                
                # Write to lab workspace
                try:
                    import re
                    import os
                    from app import SANDBOX_DIR
                    
                    u_id = getattr(g, 'user_id', None)
                    if not u_id:
                        u_cursor = conn.execute('SELECT user_id FROM chats WHERE id = ?', (chat_id,))
                        u_row = u_cursor.fetchone()
                        if u_row:
                            u_id = u_row['user_id']
                            
                    clean_uid = re.sub(r'[^a-zA-Z0-9]', '', str(u_id)) if u_id else "anon"
                    clean_cid = re.sub(r'[^a-zA-Z0-9]', '', str(chat_id))
                    workspace_name = f"lab_workspace_u{clean_uid}_c{clean_cid}"
                    lab_workspace = os.path.abspath(os.path.join(SANDBOX_DIR, workspace_name))
                    os.makedirs(lab_workspace, exist_ok=True)
                    
                    archive_path = os.path.join(lab_workspace, 'chat_history_archive.txt')
                    with open(archive_path, 'a') as f:
                        f.write(archived_messages_text)
                except Exception as lab_err:
                    logger.error("Failed to write archive to lab: %s", lab_err)
                
                placeholders = ','.join('?' * len(keep_ids))
                cursor = conn.execute(
                    f'UPDATE messages SET hidden = 1 WHERE chat_id = ? AND hidden = 0 AND id NOT IN ({placeholders})',
                    [chat_id] + keep_ids
                )
                msgs_archived = cursor.rowcount
        
        # Insert the state document as a hidden stellar message
        conn.execute(
            'INSERT INTO messages (chat_id, message_type, message_content, hidden) VALUES (?, ?, ?, ?)',
            (chat_id, 'stellar', '[COMPRESSED MEMORY STATE]\n' + state_document, 1)
        )
        
        conn.commit()
        conn.close()
        
        # Recalculate and update the token count in the DB so the UI sidebar updates immediately
        try:
            from app import count_chat_tokens
            count_chat_tokens(chat_id)
        except Exception as update_err:
            logger.error("Could not update token count after compression: %s", update_err)
            
        return f"Memory compressed successfully. Target: {target}. {tools_archived} tool calls archived. {msgs_archived} messages archived. State document preserved. Continue with your task."
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error compressing memory: {str(e)}"


def logs_and_preferences(status: str, timeout: int, write: str = "", user_id: str = "global") -> str:
    """Stores user preferences, previous errors, and resolution strategies. Memory is automatically provided to your context.
    
    Args:
        status: Status update for the user.
        write: A string detailing a preference, error, or fix to save for the future.
    """
    import sqlite3
    from app import DATABASE_NAME
    
    try:
        if not write:
            return "Error: You must provide a 'write' string to save a preference."

        write = write.strip()
        if not write:
            return "Error: 'write' string was empty."

        conn = sqlite3.connect(DATABASE_NAME)
        conn.execute("PRAGMA journal_mode=WAL;")
        conn.execute("PRAGMA busy_timeout=5000;")
        cursor = conn.cursor()
        
        # Insert new log
        cursor.execute('INSERT INTO user_logs_prefs (user_id, log_entry) VALUES (?, ?)', (user_id, write))
        
        # Keep only the last 100 entries for this user to prevent bloat
        cursor.execute('''
            DELETE FROM user_logs_prefs 
            WHERE id IN (
                SELECT id FROM user_logs_prefs 
                WHERE user_id = ? 
                ORDER BY created_at DESC 
                LIMIT -1 OFFSET 100
            )
        ''', (user_id,))
        
        conn.commit()
        conn.close()
        
        return "Successfully saved to logs/preferences. It will be available in your context for future turns."
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error accessing logs/preferences: {str(e)}"

def make_presentation(topic: str, status: str, timeout: int, num_slides: int = 10, style: str = "corporate", additional_context: str = "") -> str:
    """Generate a fully designed PowerPoint presentation where each slide is a full-bleed AI generated image containing text.
    Args:
        topic: The topic of the presentation
        status: Status update for the user.
        num_slides: Number of slides
        style: descriptive style for the presentation (e.g. "educational infographic", "minimalist executive", "dark technical documentation")
        additional_context: detailed research information to include in the slides
    """
    import asyncio
    import threading
    import json
    import os
    import uuid
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    from app import PRIMARY_API_KEY, BACKUP_API_KEYS, KEY_MANAGER, parse_quota_block_duration
    from pydantic import BaseModel, Field
    from google import genai
    from google.genai import types
    
    raw_keys = [PRIMARY_API_KEY] + [bk for bk in BACKUP_API_KEYS if bk]
    keys_to_try = [k for k in dict.fromkeys(raw_keys) if k]

    model_id = 'gemini-2.5-flash'
    active_keys = [k for k in keys_to_try if not KEY_MANAGER.is_key_blocked(k, model_id)[0]]
    if not active_keys:
        active_keys = keys_to_try
    keys_to_try = active_keys

    class Slide(BaseModel):
        """
        Pydantic model representing the content and layout description of a single slide.
        """
        title: str = Field(description="Main title for the slide.")
        summary: str = Field(description="A comprehensive, detailed summary for the slide. Provide as much informative content as necessary to make the slide educational and deep, using multiple paragraphs or extensive bullet points if needed.")
        background_description: str = Field(description="Detailed description of the visual layout: specify a multi-column infographic structure, specific diagrams (like flowcharts or state diagrams), icons, and thematic imagery that complements the dense text.")

    class PresentationPlan(BaseModel):
        """
        Pydantic model representing the full presentation plan containing a list of slides.
        """
        slides: list[Slide]

    slide_plan_prompt = (
        f"Plan {num_slides} professional infographic-style slides for a presentation on '{topic}'.\n"
        f"Style: {style}.\n"
        f"Context: {additional_context}.\n"
        "Each slide should be designed as a complete visual experience. Include specific sections, diagrams, or icons in the background description."
    )
    
    plan = None
    last_error = None
    for current_key in keys_to_try:
        try:
            client = genai.Client(api_key=current_key)
            t0 = time.time()
            resp = client.models.generate_content(
                model=model_id,
                contents=slide_plan_prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=PresentationPlan
                )
            )
            duration = time.time() - t0
            logger.info("Gemini API call completed model=%s duration_sec=%.2f purpose=make_presentation_plan", model_id, duration)
            plan = json.loads(resp.text)
            break
        except Exception as e:
            logger.error("Error in make_presentation tool error=%s", e, exc_info=True)
            error_string = str(e).lower()
            if ('429' in error_string or '403' in error_string or '503' in error_string or '500' in error_string or 'resource_exhausted' in error_string or 'quota' in error_string):
                block_duration, block_reason = parse_quota_block_duration(error_string)
                block_scope = None if ('403' in error_string or 'permission_denied' in error_string or 'invalid' in error_string) else model_id
                KEY_MANAGER.block_key(current_key, block_scope, block_duration, block_reason)
                logger.warning("Globally blocked API key hash=%s block_duration_sec=%d model=%s reason=%s error=make_presentation", hash(current_key), block_duration, block_scope, block_reason)
                last_error = e
                continue
            return f"Error: Failed to plan presentation: {str(e)}. Since presentation generation failed, please perform a web search for relevant images and place them in a nice UI for the user."
            
    if not plan:
        return f"Error: All API keys exhausted. Last error: {str(last_error)}. Since presentation generation failed, please perform a web search for relevant images and place them in a nice UI for the user."
        
    slides_data = plan.get('slides', [])
    
    async def fetch_image(slide_data):
        """
        Asynchronously generate a slide image by invoking the Gemini image generation model.

        Args:
            slide_data (dict): The slide planning details.

        Returns:
            str or None: The base64-encoded image data, or None if generation failed.
        """
        try:
            full_image_prompt = (
                f"A professional, high-resolution 16:9 presentation infographic slide. Style: '{style}'.\n"
                f"LAYOUT: {slide_data.get('background_description')}.\n"
                f"CONTENT TO RENDER DIRECTLY IN IMAGE:\n"
                f"Main Header: '{slide_data.get('title')}'\n"
                f"Body Text: '{slide_data.get('summary')}'\n"
                "INSTRUCTIONS:\n"
                "- Use professional, clean sans-serif typography.\n"
                "- Integrate the text aesthetically into a multi-column or structured infographic layout.\n"
                "- Include relevant icons, charts, or transition diagrams if mentioned in the layout description.\n"
                "- The result must be a single, complete, polished slide design. No watermarks. Legible text."
            )
            
            loop = asyncio.get_event_loop()
            def gen_sync():
                """
                Synchronous generator call wrapper run inside the event loop executor.
                """
                t0 = time.time()
                res = client.models.generate_content(
                    model='gemini-3.1-flash-image-preview',
                    contents=full_image_prompt,
                )
                duration = time.time() - t0
                logger.info("Gemini API call completed model=%s duration_sec=%.2f purpose=make_presentation_slide_image", 'gemini-3.1-flash-image-preview', duration)
                return res
            result = await loop.run_in_executor(None, gen_sync)
            if result.candidates and result.candidates[0].content.parts:
                for part in result.candidates[0].content.parts:
                    if hasattr(part, 'inline_data') and part.inline_data:
                        return part.inline_data.data
            return None
        except Exception as e:
            logger.exception("Error generating slide image: %s", e)
            return None

    async def fetch_all_images():
        """
        Fetch generated images for all slides concurrently.

        Returns:
            list: List of base64-encoded image data or None for each slide.
        """
        tasks = [fetch_image(slide) for slide in slides_data]
        return await asyncio.gather(*tasks)

    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    images = loop.run_until_complete(fetch_all_images())

    if num_slides > 0 and all(img is None for img in images):
        return "Error: Image generation model failed to generate slide images. Since presentation image generation failed, please perform a web search for relevant images and place them in a nice UI for the user."

    presentation_id = uuid.uuid4().hex[:8]
    output_dir = "outputs"
    pres_dir = os.path.join(output_dir, f"pres_{presentation_id}")
    os.makedirs(pres_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_image_urls = []
    for i, img_bytes in enumerate(images):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        slide_img_filename = f"slide_{i+1}.png"
        slide_img_path = os.path.join(pres_dir, slide_img_filename)
        
        if img_bytes:
            with open(slide_img_path, "wb") as f:
                f.write(img_bytes)
            slide_image_urls.append(f"/view/pres_{presentation_id}/{slide_img_filename}")
            
            image_stream = BytesIO(img_bytes)
            slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1)).text_frame
            title_shape.text = slides_data[i].get('title', f"Slide {i+1}")
            content_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4)).text_frame
            content_shape.text = slides_data[i].get('summary', "")
            slide_image_urls.append(None)

    pptx_filename = f"presentation_{presentation_id}.pptx"
    pptx_filepath = os.path.join(output_dir, pptx_filename)
    prs.save(pptx_filepath)
    
    result_data = {
        "presentation_id": presentation_id,
        "pptx_url": f"/download/{pptx_filename}",
        "slides": slide_image_urls,
        "topic": topic,
        "style": style,
        "num_slides": num_slides,
        "additional_context": additional_context
    }
    
    return f"PRESENTATION_DATA:{json.dumps(result_data)}"

def regenerate_presentation_slide(presentation_id: str, slide_index: int, status: str, timeout: int, topic: str = "", style: str = "", additional_context: str = "", feedback: str = "") -> str:
    """Regenerate a specific slide of an existing presentation based on feedback.
    
    Args:
        presentation_id: the ID of the presentation
        slide_index: 0-based index of the slide to regenerate
        status: Status update for the user.
        topic: original topic
        style: original style
        additional_context: original context
        feedback: specific feedback for this slide's improvement
    """
    # Inline import of genai and types to avoid startup overhead
    from google import genai
    from google.genai import types
    import os
    import json
    import asyncio
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    from app import PRIMARY_API_KEY
    from pydantic import BaseModel, Field

    from app import PRIMARY_API_KEY, BACKUP_API_KEYS, KEY_MANAGER, parse_quota_block_duration
    raw_keys = [PRIMARY_API_KEY] + [bk for bk in BACKUP_API_KEYS if bk]
    keys_to_try = [k for k in dict.fromkeys(raw_keys) if k]
    
    model_id = 'gemini-2.5-flash'
    active_keys = [k for k in keys_to_try if not KEY_MANAGER.is_key_blocked(k, model_id)[0]]
    if not active_keys:
        active_keys = keys_to_try
    keys_to_try = active_keys

    slide_plan_prompt = (
        f"We want to regenerate slide {slide_index + 1} of a presentation on '{topic}'.\n"
        f"Style: {style}.\n"
        f"Context: {additional_context}.\n"
        f"Feedback for regeneration: '{feedback}'.\n"
        "Plan the updated slide layout. Include specific sections, diagrams, or icons in the background description."
    )
    
    slide_data = None
    img_bytes = None
    last_error = None
    
    for current_key in keys_to_try:
        try:
            client = genai.Client(api_key=current_key)
            
            # Load existing slide image for reference if it exists
            existing_image_part = None
            pres_dir = os.path.join("outputs", f"pres_{presentation_id}")
            slide_path = os.path.join(pres_dir, f"slide_{slide_index + 1}.png")
            if os.path.exists(slide_path):
                with open(slide_path, "rb") as f:
                    img_data = f.read()
                    existing_image_part = {
                        'mime_type': 'image/png',
                        'data': img_data
                    }
            
            contents = [slide_plan_prompt]
            if existing_image_part:
                contents.append(types.Part.from_bytes(data=existing_image_part['data'], mime_type=existing_image_part['mime_type']))

            t0 = time.time()
            resp = client.models.generate_content(
                model=model_id,
                contents=contents,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=Slide
                )
            )
            duration = time.time() - t0
            logger.info("Gemini API call completed model=%s duration_sec=%.2f purpose=regenerate_slide_plan", model_id, duration)
            slide_data = json.loads(resp.text)
            
            full_image_prompt = (
                f"Using the attached image as a reference, modify it based on this feedback: '{feedback}'.\n"
                f"Target Infographic Layout: {slide_data.get('background_description')}.\n"
                f"CONTENT TO RENDER:\n"
                f"Main Header: '{slide_data.get('title')}'\n"
                f"Body Text: '{slide_data.get('summary')}'\n"
                "INSTRUCTIONS:\n"
                "- Maintain the professional style: '{style}'.\n"
                "- Ensure high-resolution, clean typography, and 3D infographics if requested.\n"
                "- The result must be a single, complete, polished slide design."
            )
            
            image_contents = [full_image_prompt]
            if existing_image_part:
                image_contents.append(types.Part.from_bytes(data=existing_image_part['data'], mime_type=existing_image_part['mime_type']))

            t0 = time.time()
            result = client.models.generate_content(
                model='gemini-3.1-flash-image-preview',
                contents=image_contents,
            )
            duration = time.time() - t0
            logger.info("Gemini API call completed model=%s duration_sec=%.2f purpose=regenerate_slide_image", 'gemini-3.1-flash-image-preview', duration)
            
            if result.candidates and result.candidates[0].content.parts:
                for part in result.candidates[0].content.parts:
                    if hasattr(part, 'inline_data') and part.inline_data:
                        img_bytes = part.inline_data.data
                        break
            
            break # Success!
            
        except Exception as e:
            logger.error("Error in regenerate_presentation_slide tool error=%s", e, exc_info=True)
            error_string = str(e).lower()
            if ('429' in error_string or '403' in error_string or '503' in error_string or '500' in error_string or 'resource_exhausted' in error_string or 'quota' in error_string):
                block_duration, block_reason = parse_quota_block_duration(error_string)
                block_scope = None if ('403' in error_string or 'permission_denied' in error_string or 'invalid' in error_string) else model_id
                KEY_MANAGER.block_key(current_key, block_scope, block_duration, block_reason)
                logger.warning("Globally blocked API key hash=%s block_duration_sec=%d model=%s reason=%s error=regenerate_presentation_slide", hash(current_key), block_duration, block_scope, block_reason)
                last_error = e
                continue
            return f"Error: Failed to re-plan or generate slide: {str(e)}. Since slide regeneration failed, please perform a web search for relevant images and place them in a nice UI for the user."
    if not img_bytes:
        return "Error: Failed to generate image bytes. Since slide regeneration failed, please perform a web search for relevant images and place them in a nice UI for the user."

    output_dir = "outputs"
    pres_dir = os.path.join(output_dir, f"pres_{presentation_id}")
    os.makedirs(pres_dir, exist_ok=True)
    
    slide_img_filename = f"slide_{slide_index+1}.png"
    slide_img_path = os.path.join(pres_dir, slide_img_filename)
    with open(slide_img_path, "wb") as f:
        f.write(img_bytes)
        
    # Update PPTX if it exists
    pptx_filename = f"presentation_{presentation_id}.pptx"
    pptx_filepath = os.path.join(output_dir, pptx_filename)
    if os.path.exists(pptx_filepath):
        try:
            prs = Presentation(pptx_filepath)
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                for shape in list(slide.shapes):
                    sp = shape._element
                    sp.getparent().remove(sp)
                
                image_stream = BytesIO(img_bytes)
                slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                prs.save(pptx_filepath)
        except Exception as e:
            logger.exception("Error updating PPTX pptx_filepath=%s: %s", pptx_filepath, e)

    return f"REGENERATED_SLIDE:{json.dumps({'presentation_id': presentation_id, 'slide_index': slide_index, 'url': f'/view/pres_{presentation_id}/{slide_img_filename}'})}"


def repo_control(action: str, status: str, timeout: int, app_id: str = None, project_name: str = None, files: list[str] = None, repo_url: str = None, port: int = 5000, command: str = None, env_type: str = "web") -> str:
    """Control and manage repository-based or custom-stack deployments.
    Args:
        action: "deploy", "execute", "list_history", "rename", "stop", "restart", or "snapshot"
        status: Status update for the user.
        timeout: Execution timeout in seconds.
        app_id: the Deployment ID, Project Title, or Subdomain (required for all actions except 'deploy' and 'list_history')
        project_name: Custom name for the project (used for unique subdomain in 'deploy' and 'rename')
        files: List of file paths to save into the database (required for 'snapshot')
        repo_url: URL to a git repository (optional for 'deploy')
        port: Internal port the app will listen on (default 5000, used in 'deploy')
        command: Bash command to run (required for 'execute')
        env_type: "web" or "mobile". "mobile" provisions an Android/React Native build environment.
    Returns:
        status message, history list, or command output
    """
    from app import get_db, generate_unique_subdomain, stop_and_cleanup_app_by_process_id, _redis_repo_key, redis_client, active_apps, active_apps_lock
    from flask import current_app
    session = _get_effective_session()
    import json
    import docker
    import time
    import threading
    import subprocess
    import base64
    import os

    try:
        db = get_db()

        if action == "deploy":
            if 'user_id' not in session: return "Error: Authentication required to host repos."
            
            # --- PERSISTENCE FIX: Check for existing snapshot if app_id/project_name provided ---
            existing_snapshot = None
            if app_id or project_name:
                lookup = app_id or project_name
                cursor = db.execute('SELECT files_snapshot FROM repo_history WHERE (project_name = ? OR process_id = ? OR subdomain = ?) AND user_id = ? ORDER BY id DESC LIMIT 1', (lookup, lookup, lookup, session['user_id']))
                row = cursor.fetchone()
                if row:
                    existing_snapshot = json.loads(row['files_snapshot'])
            # ----------------------------------------------------------------------------------

            process_id = str(uuid.uuid4())
            if project_name: project_title = project_name
            elif repo_url: project_title = f"Repo: {repo_url.split('/')[-1].replace('.git', '')}"
            else: project_title = "Custom Stack Project"

            subdomain = generate_unique_subdomain(project_title)
            
            logger.info("Initiating deploy action project_title=%s subdomain=%s port=%d env_type=%s process_id=%s", project_title, subdomain, port, env_type, process_id)

            # Prepare initial files snapshot
            initial_files = existing_snapshot if existing_snapshot else ({"repo": repo_url, "port": port} if repo_url else {"port": port})
            if repo_url: initial_files['repo'] = repo_url
            if port: initial_files['port'] = port

            db.execute('INSERT INTO repo_history (user_id, project_name, process_id, status, files_snapshot, subdomain) VALUES (?, ?, ?, ?, ?, ?)',
                       (session['user_id'], project_title, process_id, 'created', json.dumps(initial_files), subdomain))
            db.commit()

            try:
                client = docker.from_env()
                from app import ensure_user_network
                user_network = ensure_user_network(client, session['user_id'])
                r_key = _redis_repo_key(process_id)

                # Determine image based on env_type
                target_image = 'reactnativecommunity/react-native-android:latest' if env_type == "mobile" else 'stellar-repo-host:latest'
                
                project_dir = f"/home/stellaradmin/my_app/deployments/{process_id}"
                os.makedirs(project_dir, exist_ok=True)
                
                if existing_snapshot:
                    for fname, content in existing_snapshot.items():
                        if fname in ['repo', 'port'] or not isinstance(content, str): continue
                        fpath = os.path.join(project_dir, fname)
                        os.makedirs(os.path.dirname(fpath), exist_ok=True)
                        with open(fpath, 'wb') as f:
                            f.write(content.encode('utf-8'))

                t_run = time.time()
                container = client.containers.run(
                    image=target_image,
                    command='sleep infinity',
                    ports={f"{port}/tcp": ('0.0.0.0', 0)},                    name=f"stellar-repo-{process_id}",
                    volumes={
                        project_dir: {'bind': '/app', 'mode': 'rw'}
                    },
                    remove=False,
                    detach=True,
                    init=True,
                    network=user_network,
                    stdout=True,
                    stderr=True,
                    working_dir='/app',
                    labels={
                        "stellar_type": "repo",
                        "stellar_process_id": process_id,
                        "created_at_ts": str(time.time()),
                        "repo_app_id": process_id,
                        "subdomain": subdomain
                    }
                )
                logger.info("Repo container created process_id=%s image=%s duration_sec=%.2f", process_id, target_image, time.time() - t_run)
                time.sleep(2)
                container.reload()
                host_port = container.attrs['NetworkSettings']['Ports'][f"{port}/tcp"][0]['HostPort']
                redis_client.hset(r_key, mapping={"container_id": container.id, "status": "running", "process_id": process_id, "host_port": str(host_port), "files": json.dumps(initial_files)})
                with active_apps_lock: active_apps[process_id] = {"container_id": container.id, "port": host_port, "status": "running"}
                db.execute("UPDATE repo_history SET status = 'running', deployment_url = ? WHERE process_id = ?", (f"https://{subdomain}.stellarai.live/", process_id))
                db.commit()
                
                if repo_url and not os.listdir(project_dir):
                    clone_res = container.exec_run(f"git clone {repo_url} .")
                    if clone_res.exit_code != 0: return f"Git clone failed: {clone_res.output.decode()}"
                
                restored_count = len(existing_snapshot) if existing_snapshot else 0
                
                public_url = f"https://{subdomain}.stellarai.live/"
                restore_msg = f" (Restored {restored_count} files from snapshot)" if restored_count > 0 else ""
                return f"Container provisioned for '{project_title}'! ID: `{process_id}`. Live URL: {public_url}{restore_msg} Use action='execute' to build and start the app. CRITICAL: Ensure your app listens on 0.0.0.0 and port {port}."
            except Exception as e:
                logger.exception("Error caught: %s", e)
                db.execute("UPDATE repo_history SET status = 'failed' WHERE process_id = ?", (process_id,))
                db.commit()
                return f"Error provisioning container: {str(e)}"

        if action == "execute":
            if not app_id or not command: return "Error: 'app_id' and 'command' are required for execute."
            cursor = db.execute('SELECT process_id, files_snapshot FROM repo_history WHERE (project_name = ? OR process_id = ? OR subdomain = ?) AND user_id = ? ORDER BY id DESC LIMIT 1', (app_id, app_id, app_id, session.get('user_id')))
            row = cursor.fetchone()
            if not row: return f"Error: Deployment '{app_id}' not found."
            p_id = row['process_id']
            snapshot = json.loads(row['files_snapshot'])
            port = snapshot.get('port', 5000)
            
            logger.info("Executing command in repo container app_id=%s process_id=%s command=%s", app_id, p_id, command)
            t_exec = time.time()
            try:
                client = docker.from_env()
                container = client.containers.get(f"stellar-repo-{p_id}")
                wrapped_cmd = f"bash -c {subprocess.list2cmdline([command])}"
                
                # --- OCI RUNTIME FIX: Avoid explicit workdir to prevent namespace errors ---
                exec_result = container.exec_run(wrapped_cmd, demux=False)
                output = exec_result.output.decode('utf-8', 'replace')
                duration = time.time() - t_exec
                logger.info("Command execution completed app_id=%s process_id=%s duration_sec=%.3f exit_code=%d", app_id, p_id, duration, exec_result.exit_code)

                # --- OCI RUNTIME ERROR RECOVERY ---
                if exec_result.exit_code == 128 and ("mount namespace root" in output or "container breakout" in output or "exec failed" in output):
                    logger.warning(f"Detected broken container mount for stellar-repo-{p_id}. Restarting...")
                    try:
                        # Attempt a quick stop and restart to refresh mounts
                        container.stop(timeout=2)
                        container.remove(force=True)
                    except Exception as clean_err:
                        logger.warning("Failed to clean up container on OCI error process_id=%s: %s", p_id, clean_err)
                    
                    restart_res = repo_control(action="restart", status="Restarting for recovery...", app_id=p_id)
                    if "Error" in restart_res: return restart_res
                    
                    # Prevent infinite recursion
                    if kwargs.get('_retry_count', 0) < 2:
                        kwargs['_retry_count'] = kwargs.get('_retry_count', 0) + 1
                        return repo_control(action, status, timeout, app_id, project_name, files, repo_url, port, command, env_type, **kwargs)
                # ----------------------------------
                
                # Enhanced health check if a start-like command is detected
                if any(kw in command.lower() for kw in ["npm start", "python", "node", "serve", "go run", "npm run dev"]):
                    time.sleep(5)
                    container.reload()
                    if container.status != 'running':
                        return f"Command executed, but container stopped. Output:\n{output}"
                    
                    # Verify port accessibility and status
                    status_code = 0
                    for _ in range(5):
                        time.sleep(2)
                        res = container.exec_run(f"curl -s -o /dev/null -w '%{{http_code}}' http://localhost:{port}/")
                        if res.exit_code == 0:
                            status_code = int(res.output.decode().strip())
                            if 0 < status_code < 500:
                                break
                    
                    if 0 < status_code < 500:
                        output += f"\n\n✨ Server is READY (Status {status_code}) and listening on 0.0.0.0:{port}!"
                    elif status_code >= 500:
                        output += f"\n\n❌ Server responded with ERROR {status_code} on port {port}. Check your logs!"
                    else:
                        output += f"\n\n⚠️ Command run, but server is NOT responding on port {port}. Ensure it listens on 0.0.0.0."

                if exec_result.exit_code != 0: return f"Command failed (exit {exec_result.exit_code}).\nOutput:\n{output}"
                return output if output else "Command executed successfully with no output."
            except Exception as e:
                logger.exception("Error caught: %s", e)
                return f"Error executing in repo: {str(e)}"

        if action == "list_history":
            if 'user_id' not in session: return "Error: Authentication required."
            cursor = db.execute('SELECT project_name, process_id, status, deployment_url, subdomain, created_at FROM repo_history WHERE user_id = ? ORDER BY created_at DESC', (session['user_id'],))
            history = cursor.fetchall()
            if not history: return "You have no past deployments."
            res = "### Your Deployment History:\n"
            for row in history:
                url = row['deployment_url'] or (f"https://{row['subdomain']}.stellarai.live/" if row['subdomain'] else f"https://{row['process_id']}.stellarai.live/")
                url_str = f" - [Visit App]({url})" if row['status'] == 'running' else ""
                res += f"- **{row['project_name']}** (ID: `{row['process_id']}`) - Status: {row['status']} - Created: {row['created_at']}{url_str}\n"
            return res

        if action == "rename":
            if not app_id or not project_name:
                return "Error: Both 'app_id' (current project) and 'project_name' (new name) are required for rename."
            
            cursor = db.execute('SELECT process_id, project_name FROM repo_history WHERE (project_name = ? OR process_id = ? OR subdomain = ?) AND user_id = ? ORDER BY id DESC LIMIT 1', (app_id, app_id, app_id, session.get('user_id')))
            row = cursor.fetchone()
            if not row: return f"Error: Deployment '{app_id}' not found."
            
            actual_id = row['process_id']
            new_subdomain = generate_unique_subdomain(project_name)
            new_url = f"https://{new_subdomain}.stellarai.live/"
            
            logger.info("Renaming deployment app_id=%s process_id=%s to new_name=%s new_subdomain=%s", app_id, actual_id, project_name, new_subdomain)
            
            db.execute("UPDATE repo_history SET project_name = ?, subdomain = ?, deployment_url = ? WHERE process_id = ?", (project_name, new_subdomain, new_url, actual_id))
            db.commit()
            return f"Deployment renamed to '{project_name}'! New URL: {new_url}"

        def _perform_snapshot(p_id, container_name, current_snapshot):
            """
            Snapshot the container's file system state back to the database.

            Args:
                p_id (str): The process ID of the deployment.
                container_name (str): The name of the Docker container.
                current_snapshot (dict): The current tracked file mapping in the snapshot.
            """
            try:
                client = docker.from_env()
                container = client.containers.get(container_name)
                # We identify files to snapshot by looking at the current snapshot keys
                # This ensures we don't accidentally snapshot huge binary folders like node_modules
                # unless they were part of the project's tracked files.
                # However, for repo, we should also look for known code files.
                tracked_files = list(current_snapshot.keys())
                
                # Scan for any code files in /app to ensure new files are captured
                res = container.exec_run("find . -maxdepth 3 -not -path '*/.*' -not -path './node_modules/*' -not -path './venv/*' -type f")
                if res.exit_code == 0:
                    found_files = res.output.decode('utf-8', 'replace').strip().split('\n')
                    for f in found_files:
                        f = f.lstrip('./')
                        if f and f not in tracked_files:
                            tracked_files.append(f)
                
                count = 0
                for file_path in tracked_files:
                    # Remove leading /app/ if present
                    clean_path = file_path.replace('/app/', '').lstrip('/')
                    if clean_path in ['repo', 'port']: continue
                    
                    res = container.exec_run(f"cat {clean_path}")
                    if res.exit_code == 0:
                        current_snapshot[clean_path] = res.output.decode('utf-8', 'replace')
                        count += 1
                
                db.execute("UPDATE repo_history SET files_snapshot = ? WHERE process_id = ?", (json.dumps(current_snapshot), p_id))
                db.commit()
                return count
            except Exception as e:
                logger.exception("Auto-snapshot failed for process_id=%s: %s", p_id, e)
                return 0

        if action == "stop":
            if not app_id: return "Error: app_id is required to stop a deployment."
            cursor = db.execute('SELECT process_id, files_snapshot FROM repo_history WHERE (project_name = ? OR process_id = ? OR subdomain = ?) AND user_id = ? ORDER BY id DESC LIMIT 1', (app_id, app_id, app_id, session.get('user_id')))
            row = cursor.fetchone()
            if not row: return f"Error: Deployment '{app_id}' not found."
            
            p_id = row['process_id']
            current_snapshot = json.loads(row['files_snapshot'])
            
            current_snapshot = json.loads(row['files_snapshot'])
            
            logger.info("Stopping deployment app_id=%s process_id=%s", app_id, p_id)
            
            stop_and_cleanup_app_by_process_id(p_id, app_type='repo')
            db.execute("UPDATE repo_history SET status = 'stopped' WHERE process_id = ?", (p_id,))
            db.commit()
            
            return f"Deployment '{app_id}' has been stopped. Files are safely persisted on the host filesystem."

        if action == "restart":
            if not app_id: return "Error: app_id is required to restart a deployment."
            
            cursor = db.execute('SELECT process_id, project_name, files_snapshot, subdomain FROM repo_history WHERE (project_name = ? OR process_id = ? OR subdomain = ?) AND user_id = ? ORDER BY id DESC LIMIT 1', (app_id, app_id, app_id, session.get('user_id')))
            row = cursor.fetchone()
            if not row: return f"Error: Deployment '{app_id}' not found."
            
            process_id = row['process_id']
            project_title = row['project_name']
            subdomain = row['subdomain']
            current_snapshot = json.loads(row['files_snapshot'])
            
            current_snapshot = json.loads(row['files_snapshot'])
            
            logger.info("Restarting deployment app_id=%s process_id=%s subdomain=%s", app_id, process_id, subdomain)
            t_restart = time.time()

            # Reload updated snapshot
            cursor = db.execute('SELECT files_snapshot FROM repo_history WHERE process_id = ?', (process_id,))
            row = cursor.fetchone()
            snapshot = json.loads(row['files_snapshot'])
            
            repo_url = snapshot.get('repo')
            port = snapshot.get('port', 5000)
            

            # Generic Repo/Custom restart
            stop_and_cleanup_app_by_process_id(process_id, app_type='repo')
            
            client = docker.from_env()
            from app import ensure_user_network
            user_network = ensure_user_network(client, session['user_id'])
            r_key = _redis_repo_key(process_id)
            
            project_dir = f"/home/stellaradmin/my_app/deployments/{process_id}"
            os.makedirs(project_dir, exist_ok=True)
            
            if not os.listdir(project_dir) and snapshot:
                for fname, content in snapshot.items():
                    if fname in ['repo', 'port'] or not isinstance(content, str): continue
                    fpath = os.path.join(project_dir, fname)
                    os.makedirs(os.path.dirname(fpath), exist_ok=True)
                    with open(fpath, 'wb') as f:
                        f.write(content.encode('utf-8'))

            t_run = time.time()
            container = client.containers.run(
                image='stellar-repo-host:latest',
                command='sleep infinity',
                ports={f"{port}/tcp": ('0.0.0.0', 0)},
                volumes={
                    project_dir: {'bind': '/app', 'mode': 'rw'}
                },
                name=f"stellar-repo-{process_id}",
                detach=True,
                init=True,
                working_dir='/app',
                network=user_network,
                labels={
                    "stellar_type": "repo",
                    "stellar_process_id": process_id,
                    "created_at_ts": str(time.time()),
                    "repo_app_id": process_id,
                    "subdomain": subdomain
                }
            )
            logger.info("Repo container restarted process_id=%s image=%s duration_sec=%.2f", process_id, 'stellar-repo-host:latest', time.time() - t_run)

            time.sleep(2)
            container.reload()
            host_port = container.attrs['NetworkSettings']['Ports'][f"{port}/tcp"][0]['HostPort']

            redis_client.hset(r_key, mapping={"container_id": container.id, "status": "running", "process_id": process_id, "host_port": str(host_port), "files": json.dumps(snapshot)})
            with active_apps_lock: active_apps[process_id] = {"container_id": container.id, "port": host_port, "status": "running"}

            if repo_url and not os.listdir(project_dir):
                container.exec_run(f"git clone {repo_url} .")

            public_url = f"https://{subdomain}.stellarai.live/" if subdomain else f"https://{process_id}.stellarai.live/"
            
            # Re-run health check loop
            status_code = 0
            for _ in range(30):
                time.sleep(1)
                try:
                    res = container.exec_run(f"curl -s -o /dev/null -w '%{{http_code}}' http://localhost:{port}/")
                    if res.exit_code == 0:
                        status_code = int(res.output.decode().strip())
                        if 0 < status_code < 500:
                            break
                        elif status_code >= 500:
                            break # Fail fast on server error
                except Exception as health_err:
                    logger.debug("Container restart health check failed on port %d: %s", port, health_err)
            
            logger.info("Restart completed app_id=%s process_id=%s status_code=%d duration_sec=%.3f", app_id, process_id, status_code, time.time() - t_restart)

            if 0 < status_code < 500:
                ready_msg = f"and server is READY (Status {status_code})!"
            elif status_code >= 500:
                ready_msg = f"but server returned ERROR {status_code}. Check your logs!"
            else:
                ready_msg = "but server is not responding yet. Use repo_control action='execute' to start your app."
                
            return f"Deployment '{project_title}' restarted with latest snapshotted edits {ready_msg} ID: `{process_id}`. Live URL: {public_url}"

        if action == "snapshot":
            if not app_id or not files:
                return "Error: app_id and a list of 'files' paths are required to snapshot manual edits."
            
            return f"Successfully snapshotted {len(files)} files natively to the host filesystem for project '{app_id}' (DB snapshots are deprecated)."

        return f"Error: Unknown action '{action}'."
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error in repo_control: {str(e)}"

# Remove host_repo and repo_execute functions and their available_tools entries



def lab_execute(command: str, status: str, timeout: int) -> str:
    """Executes a bash command in a persistent, isolated Docker sandbox.
    Args:
        command: The bash command to run.
        status: Status update for the user.
        timeout: Execution timeout in seconds.
    """
    import subprocess
    import docker
    from app import SANDBOX_DIR, UPLOAD_FOLDER
    import os
    import re
    import shutil
    session = _get_effective_session()

    # Deterministic Context Detection
    u_id = None
    c_id = 'default'
    s_id = 'no_session'

    if has_request_context():
        u_id = session.get('user_id')
        c_id = session.get('current_chat_id', 'default')
        s_id = getattr(session, 'sid', 'no_session')

    if not u_id: u_id = getattr(g, 'user_id', None)
    if c_id == 'default': c_id = getattr(g, 'chat_id', 'default')
    if s_id == 'no_session': s_id = getattr(g, 'session_id', 'no_session')

    # Sanitize for Docker/FS safety
    clean_uid = re.sub(r'[^a-zA-Z0-9]', '', str(u_id)) if u_id else "anon"
    clean_cid = re.sub(r'[^a-zA-Z0-9]', '', str(c_id))

    # Workspace naming now linked to User+Chat for background continuity
    container_name = f"stellar-lab-u{clean_uid}-c{clean_cid}"
    workspace_name = f"lab_workspace_u{clean_uid}_c{clean_cid}"

    image_name = "stellar-lab-core:latest"    
    try:
        client = docker.from_env()
        from app import ensure_user_network
        user_network = ensure_user_network(client, u_id)
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error: Docker client not available: {str(e)}"
        
    # Create a private workspace for the lab
    lab_workspace = os.path.abspath(os.path.join(SANDBOX_DIR, workspace_name))
    os.makedirs(lab_workspace, exist_ok=True)

    # --- DETERMINISTIC FIX: Auto-sync uploaded files to the lab workspace ---
    try:
        if c_id and c_id != 'no_session':
            local_dir = os.path.join(UPLOAD_FOLDER, str(c_id))
            if os.path.exists(local_dir):
                for f in os.listdir(local_dir):
                    src = os.path.join(local_dir, f)
                    dst = os.path.join(lab_workspace, f)
                    if os.path.isfile(src):
                        shutil.copy2(src, dst)
    except Exception as sync_e:
        logger.exception("Error caught: %s", sync_e)
        pass # Silent fail if permissions or path issues
    # ------------------------------------------------------------------------

    # Ensure sandbox container is running
    container = None
    try:
        container = client.containers.get(container_name)
        if container.status != 'running':
            t_start = time.time()
            container.start()
            logger.info("Lab container started name=%s duration_sec=%.2f", container_name, time.time() - t_start)
    except docker.errors.NotFound:
        try:
            t_run = time.time()
            container = client.containers.run(
                image_name,
                name=container_name,
                detach=True,
                tty=True,
                init=True,
                working_dir='/lab',
                volumes={lab_workspace: {'bind': '/lab', 'mode': 'rw'}},
                restart_policy={"Name": "unless-stopped"},
                network=user_network,
                labels={
                    "stellar_type": "lab",
                    "stellar_user_id": str(u_id),
                    "stellar_chat_id": str(c_id),
                    "created_at_ts": str(time.time())
                }
            )
            logger.info("Lab container run completed name=%s duration_sec=%.2f", container_name, time.time() - t_run)
        except Exception as e:
            logger.exception("Error caught: %s", e)
            return f"Failed to start Lab sandbox: {str(e)}"
    
    # Execute the command
    logger.info("Executing command in Lab sandbox container_name=%s command=%s", container_name, command)
    t_cmd = time.time()
    try:
        # Wrap command to capture stdout and stderr together and handle errors gracefully
        wrapped_cmd = f"bash -c {subprocess.list2cmdline([command])}"
        
        # --- OCI RUNTIME FIX: Avoid explicit workdir if possible to prevent namespace errors ---
        # The container is already started with working_dir='/lab' bound to the workspace.
        exec_result = container.exec_run(
            wrapped_cmd,
            demux=False # Get combined stdout/stderr
        )
        
        output = exec_result.output.decode('utf-8', 'replace')
        exit_code = exec_result.exit_code
        duration = time.time() - t_cmd
        logger.info("Lab command executed container_name=%s duration_sec=%.3f exit_code=%d", container_name, duration, exit_code)
        
        # --- OCI RUNTIME ERROR RECOVERY ---
        # Detects 'current working directory is outside of container mount namespace root'
        # and other fatal OCI/container breakout errors that require a fresh sandbox.
        is_oci_error = (exit_code == 128 and ("mount namespace root" in output or "container breakout" in output or "exec failed" in output))
        
        if is_oci_error:
            # Prevent infinite recursion using Flask's g to track retries for this chat/container
            retry_count = 0
            from flask import has_app_context
            if has_app_context():
                retry_key = f"lab_retry_{container_name}"
                retry_count = getattr(g, retry_key, 0)
                setattr(g, retry_key, retry_count + 1)
            
            if retry_count < 2:
                logger.warning(f"Detected broken container mount for {container_name} (retry {retry_count}). Recreating...")
                try:
                    container.stop(timeout=2)
                    container.remove(force=True)
                except:
                    pass
                time.sleep(1) # Grace period for Docker to release resources
                return lab_execute(command, status, timeout)
            else:
                return f"Critical OCI Runtime Error (Persistent after recreation): {output}"
        # ----------------------------------
        
        if exit_code != 0:
            return f"Command failed with exit code {exit_code}.\nOutput:\n{output}"
        
        return output if output else "Command executed successfully with no output."
        
    except Exception as e:
        error_msg = str(e)
        if "mount namespace root" in error_msg or "container breakout" in error_msg:
            logger.warning(f"Caught OCI error exception for {container_name}. Attempting recovery...")
            try:
                container.stop(timeout=2)
                container.remove(force=True)
            except:
                pass
            time.sleep(1)
            return lab_execute(command, status, timeout)

        logger.exception("Error caught during lab_execute: %s", e)
        return f"Error executing command in Lab: {str(e)}"

def obtain_talent(talent_names: list[str], status: str, timeout: int, search_query: str = None) -> str:
    """
    Load specialized mandates (talents) from the database to acquire detailed instructions, rules, and best practices for specific roles (e.g., frontend_design, generative_ai).
    You can pass exact talent_names OR a search_query to search for keywords across all talents.
    This tool is protected from output truncation.

    Args:
        talent_names (list[str]): The names of the talents/mandates to obtain (e.g., ['generative_ai', 'frontend_design']). Leave empty if using search_query.
        status (str): The status to report to the user.
        timeout (int): The maximum time to wait for the command to complete.
        search_query (str): Optional. A keyword to search across all templates/talents (e.g., 'cars', 'ecommerce').
    """
    import sqlite3
    import os
    db_path = os.path.join(os.path.dirname(__file__), 'stellar_local.db')
    try:
        conn = sqlite3.connect(db_path)
        conn.execute("PRAGMA journal_mode=WAL;")
        conn.execute("PRAGMA busy_timeout=5000;")
        c = conn.cursor()
        
        results = []
        if search_query:
            query = f"%{search_query}%"
            c.execute("SELECT talent_name, mandate_text FROM talents WHERE mandate_text LIKE ? OR talent_name LIKE ?", (query, query))
            rows = c.fetchall()
            if rows:
                results.append(f"--- SEARCH RESULTS FOR '{search_query}' ---")
                for row in rows:
                    results.append(f"Talent Name: {row[0]}\nPreview:\n{row[1][:1500]}...\n")
            else:
                results.append(f"No talents found matching '{search_query}'.")

        if talent_names:
            for name in talent_names:
                c.execute("SELECT mandate_text FROM talents WHERE talent_name = ?", (name.lower(),))
                row = c.fetchone()
                if row:
                    results.append(f"--- TALENT ACQUIRED: {name.upper()} ---\n{row[0]}")
                else:
                    results.append(f"Error: Talent '{name}' not found in database.")
                
        conn.close()
        return "\n\n".join(results)
    except Exception as e:
        return f"Error obtaining talent: {str(e)}"

def read_tool_output(output_id: int, status: str, timeout: int, keyword: str = None, start_line: int = 0, max_lines: int = 100) -> str:
    """Reads a specific slice of a past tool's output from the database.
    Use this when history says [Output truncated] to retrieve data without context overflow.
    
    Keyword Search: If a 'keyword' is provided, the tool returns only lines containing that keyword.
    Pagination: 'start_line' acts as an offset (either for raw lines or for keyword matches).
    'max_lines' limits the number of lines returned in one call.
    
    Args:
        output_id: The ID of the tool execution to read.
        status: Status update for the user.
        timeout: Execution timeout in seconds.
        keyword: Optional string to search for.
        start_line: The line number (or match index) to start from (0-indexed).
        max_lines: The maximum number of lines to return.
    """
    import sqlite3
    try:
        from app import DATABASE_NAME
        conn = sqlite3.connect(DATABASE_NAME)
        conn.execute("PRAGMA journal_mode=WAL;")
        conn.execute("PRAGMA busy_timeout=5000;")
        conn.row_factory = sqlite3.Row
        cursor = conn.execute('SELECT result FROM tool_calls WHERE id = ?', (output_id,))
        row = cursor.fetchone()
        conn.close()

        if not row:
            return f"Error: No tool output found with ID {output_id}."

        res_str = str(row['result'])
        lines = res_str.split('\n')
        
        if keyword:
            # Filter lines by keyword and return them with line numbers
            matching_lines = []
            for i, line in enumerate(lines):
                if keyword.lower() in line.lower():
                    matching_lines.append(f"Line {i}: {line}")
            
            if not matching_lines:
                return f"No occurrences of '{keyword}' found in tool output {output_id}."
            
            total_matches = len(matching_lines)
            end_line = min(start_line + max_lines, total_matches)
            sliced_matches = matching_lines[start_line:end_line]
            
            output = f"--- Tool Output ID: {output_id} (Keyword: '{keyword}', matches {start_line} to {end_line-1} of {total_matches}) ---\n"
            output += '\n'.join(sliced_matches)
            if end_line < total_matches:
                output += f"\n--- (Too many matches. {total_matches - end_line} more lines containing '{keyword}'. Read from line {end_line} to see more) ---"
            return output

        total_lines = len(lines)
        if start_line >= total_lines:
            return f"Error: start_line {start_line} is beyond the total lines ({total_lines})."

        end_line = min(start_line + max_lines, total_lines)
        sliced_lines = lines[start_line:end_line]

        output = f"--- Tool Output ID: {output_id} (Lines {start_line} to {end_line-1} of {total_lines}) ---\n"
        output += '\n'.join(sliced_lines)
        if end_line < total_lines:
            output += f"\n--- (Output truncated. {total_lines - end_line} lines remaining. Read from line {end_line} to see more) ---"
        
        return output
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error reading tool output: {str(e)}"

def analyze_youtube_video(query: str, status: str, timeout: int, action: str = "analyze", video_url: str = None, start_time: str = None, end_time: str = None, fps: int = 1, max_results: int = 5, model_id: str = "gemma-4-31b-it") -> str:
    """Analyzes a specific YouTube video or searches for the best video based on a query.
    Args:
        query: What you want to find/analyze. Mandatory for both actions.
        status: Status update for the user.
        action: 'analyze' (default) to interrogate a video's content, or 'search' to find videos matching the query.
        video_url: The full YouTube URL (required for 'analyze').
        start_time: Optional start offset for 'analyze' (e.g., '1m10s', '60s').
        end_time: Optional end offset for 'analyze' (e.g., '2m30s', '120s').
        fps: Frames per second to sample from the video for 'analyze' (default 1).
        max_results: Number of search results to return (default 5, max 50).
        model_id: (Internal) The model to use for analysis.
    """
    from app import PRIMARY_API_KEY, YOUTUBE_API_KEY
    import requests
    
    if action == "search":
        if not YOUTUBE_API_KEY:
            return "Error: YouTube API key is not configured in the backend."
        try:
            # Enforce max 50 limit
            limit = min(int(max_results), 50)
            
            # Step 1: Search for Video IDs
            search_url = "https://www.googleapis.com/youtube/v3/search"
            search_params = {
                "part": "snippet",
                "q": query,
                "maxResults": limit,
                "type": "video",
                "key": YOUTUBE_API_KEY
            }
            # Bolt - Stability Optimization: Add timeout=15 to prevent threads hanging indefinitely on external API calls
            search_response = requests.get(search_url, params=search_params, timeout=15).json()
            if "error" in search_response:
                return f"YouTube Search API Error: {json.dumps(search_response['error'])}"
            
            items = search_response.get("items", [])
            video_ids = [item["id"]["videoId"] for item in items]
            if not video_ids:
                return "No YouTube videos found for the given query."

            # Step 2: Get detailed stats (Views, Likes, Duration, and Full Description)
            stats_url = "https://www.googleapis.com/youtube/v3/videos"
            stats_params = {
                "part": "snippet,statistics,contentDetails",
                "id": ",".join(video_ids),
                "key": YOUTUBE_API_KEY
            }
            # Bolt - Stability Optimization: Add timeout=15 to prevent threads hanging indefinitely on external API calls
            stats_response = requests.get(stats_url, params=stats_params, timeout=15).json()
            
            enriched_results = []
            for stat_item in stats_response.get("items", []):
                v_id = stat_item["id"]
                snippet = stat_item.get("snippet", {})
                
                enriched_results.append({
                    "title": snippet.get("title"),
                    "channelTitle": snippet.get("channelTitle"),
                    "description": snippet.get("description"),
                    "video_id": v_id,
                    "viewCount": int(stat_item.get("statistics", {}).get("viewCount", 0)),
                    "likeCount": int(stat_item.get("statistics", {}).get("likeCount", 0)),
                    "duration": stat_item.get("contentDetails", {}).get("duration", "N/A"),
                    "url": f"https://www.youtube.com/watch?v={v_id}"
                })
            
            # Sort by viewCount descending
            enriched_results.sort(key=lambda x: x["viewCount"], reverse=True)
            return json.dumps(enriched_results, indent=2)
        except Exception as e:
            logger.exception("Error caught: %s", e)
            return f"Error during YouTube search: {str(e)}"

    # Default 'analyze' logic
    if not video_url:
        return "Error: 'video_url' is required for the 'analyze' action."

    # Inline import of genai and types to avoid startup overhead
    from google import genai
    from google.genai import types
    from app import PRIMARY_API_KEY, BACKUP_API_KEYS, KEY_MANAGER, parse_quota_block_duration
    raw_keys = [PRIMARY_API_KEY] + [bk for bk in BACKUP_API_KEYS if bk]
    keys_to_try = [k for k in dict.fromkeys(raw_keys) if k]
    
    active_keys = [k for k in keys_to_try if not KEY_MANAGER.is_key_blocked(k, model_id)[0]]
    if not active_keys:
        active_keys = keys_to_try
    keys_to_try = active_keys
    
    part = types.Part(
        file_data=types.FileData(
            file_uri=video_url,
            mime_type="video/*",
        )
    )
    
    video_metadata = {}
    if start_time: video_metadata['start_offset'] = start_time
    if end_time: video_metadata['end_offset'] = end_time
    if fps: video_metadata['fps'] = int(fps)
    
    if video_metadata:
        part.video_metadata = types.VideoMetadata(**video_metadata)
        
    contents = [
        types.Content(
            role="user",
            parts=[part, types.Part.from_text(text=query)]
        )
    ]
    
    last_error = None
    for current_key in keys_to_try:
        try:
            client = genai.Client(api_key=current_key, http_options={'api_version': 'v1beta'})
            t0 = time.time()
            response = client.models.generate_content(
                model=model_id,
                contents=contents
            )
            duration = time.time() - t0
            logger.info("Gemini API call completed model=%s duration_sec=%.2f purpose=analyze_youtube_video", model_id, duration)
            return response.text if response.text else "The model returned an empty response for the video analysis."
        except Exception as e:
            logger.error("Error in analyze_youtube_video tool error=%s", e, exc_info=True)
            error_string = str(e).lower()
            if ('429' in error_string or '403' in error_string or '503' in error_string or '500' in error_string or 'resource_exhausted' in error_string or 'quota' in error_string):
                block_duration, block_reason = parse_quota_block_duration(error_string)
                block_scope = None if ('403' in error_string or 'permission_denied' in error_string or 'invalid' in error_string) else model_id
                KEY_MANAGER.block_key(current_key, block_scope, block_duration, block_reason)
                logger.warning("Globally blocked API key hash=%s block_duration_sec=%d model=%s reason=%s error=analyze_youtube_video", hash(current_key), block_duration, block_scope, block_reason)
                last_error = e
                continue
            return f"Error analyzing YouTube video: {str(e)}"

    return f"Error: All API keys exhausted. Last error: {str(last_error)}"

def manage_files(action: str, status: str, timeout: int, file_name: str = None, target_env: str = "lab", source_env: str = "chat") -> str:
    """
    Manage user-uploaded files or export code out of execution environments.
    Args:
        action: 'read' (list uploads), 'move' (transfer file/folder), 'project' (export to user).
        status: Status update for the user.
        file_name: The name of the file or folder to move or project.
        target_env: 'lab' or process_id for repo.
        source_env: 'chat' (uploads), 'lab', or process_id for repo.
    """
    session = _get_effective_session()
    from app import app, UPLOAD_FOLDER, get_db
    import os
    import docker
    import base64
    import uuid
    import tarfile
    import io
    import re

    try:
        client = docker.from_env()
    except Exception as e:
        logger.exception("Error caught: %s", e)
        return f"Error: Docker client not available: {str(e)}"

    from flask import g, has_request_context
    u_id = None
    c_id = 'default'
    s_id = 'no_session'

    if has_request_context():
        u_id = session.get('user_id')
        c_id = session.get('current_chat_id', 'default')
        s_id = getattr(session, 'sid', 'no_session')

    if not u_id: u_id = getattr(g, 'user_id', None)
    if c_id == 'default': c_id = getattr(g, 'chat_id', 'default')
    if s_id == 'no_session': s_id = getattr(g, 'session_id', 'no_session')

    clean_uid = re.sub(r'[^a-zA-Z0-9]', '', str(u_id)) if u_id else "anon"
    clean_cid = re.sub(r'[^a-zA-Z0-9]', '', str(c_id))
    
    dynamic_lab_container = f"stellar-lab-u{clean_uid}-c{clean_cid}"
    context_id = str(c_id) if c_id and c_id != 'default' else str(s_id)

    def resolve_env_id(env_id):
        """
        Resolve an environment identifier (which can be a project name, subdomain, or process ID)
        to the canonical process ID.

        Args:
            env_id (str): The environment identifier to resolve.

        Returns:
            str: The canonical process ID, or the original identifier if unresolved.
        """
        if not env_id or env_id in ("lab", "chat"): return env_id
        try:
            db = get_db()
            cursor = db.execute('SELECT process_id FROM repo_history WHERE (project_name = ? OR process_id = ? OR subdomain = ?) AND user_id = ? ORDER BY id DESC LIMIT 1', (env_id, env_id, env_id, u_id))
            row = cursor.fetchone()
            if row: return row[0]
        except: pass
        return env_id

    target_env = resolve_env_id(target_env)
    source_env = resolve_env_id(source_env)

    def validate_env_ownership(env_id):
        """
        Validate that the current user owns/has permission to access the specified environment.

        Args:
            env_id (str): The process ID or environment name to validate.

        Returns:
            bool: True if ownership is valid; False otherwise.
        """
        if env_id in ("lab", "chat"): return True
        try:
            db = get_db()
            cursor = db.execute('SELECT 1 FROM repo_history WHERE process_id = ? AND user_id = ?', (env_id, u_id))
            if cursor.fetchone(): return True
            if has_request_context():
                if env_id == session.get('last_run_code_process_id'): return True
        except: pass
        return False

    if not validate_env_ownership(target_env):
        return f"Error: Unauthorized access to target environment '{target_env}'."
    if not validate_env_ownership(source_env):
        return f"Error: Unauthorized access to source environment '{source_env}'."

    if target_env == "lab" or source_env == "lab":
        try:
            client.containers.get(dynamic_lab_container)
        except Exception:
            logger.exception("Error caught.")
            lab_execute("echo 'init'", "Initializing Lab...", 60)

    def get_container_name(env_id):
        """
        Get the Docker container name for the given environment ID.

        Args:
            env_id (str): The environment identifier (e.g. 'lab' or a process ID).

        Returns:
            str: The Docker container name.
        """
        if env_id == "lab": return dynamic_lab_container
        return f"stellar-repo-{env_id}"

    local_uploads = os.path.join(UPLOAD_FOLDER, context_id)

    if action == "read":
        if not os.path.exists(local_uploads): return "No uploads found."
        files = os.listdir(local_uploads)
        return "Files in chat context:\n" + "\n".join([f"- {f}" for f in files]) if files else "No uploads found."

    elif action == "move":
        if not file_name: return "Error: 'file_name' required."
        
        # Determine source and target containers/paths
        target_name = get_container_name(target_env)
        target_dir = "/lab" if target_env == "lab" else "/app"
        
        try:
            target_cont = client.containers.get(target_name)
            
            # --- SOURCE: CHAT (LOCAL UPLOADS) ---
            if source_env == "chat":
                src_path = os.path.join(local_uploads, file_name)
                if not os.path.exists(src_path): return f"File '{file_name}' not found in uploads."
                
                # Security: Restrict host moves to only UPLOAD_FOLDER and OUTPUTS
                allowed_prefixes = [os.path.abspath(UPLOAD_FOLDER), os.path.abspath("outputs")]
                if not any(os.path.abspath(src_path).startswith(p) for p in allowed_prefixes):
                    return "Security Error: Host move restricted to specific directories."

                tar_stream = io.BytesIO()
                with tarfile.open(fileobj=tar_stream, mode='w') as tar:
                    tar.add(src_path, arcname=os.path.basename(src_path))
                tar_stream.seek(0)
                if target_cont.put_archive(target_dir, tar_stream):
                    return f"Moved '{file_name}' from chat to {target_env}."
            
            # --- SOURCE: ANOTHER CONTAINER ---
            else:
                source_name = get_container_name(source_env)
                source_dir = "/lab" if source_env == "lab" else "/app"
                src_cont = client.containers.get(source_name)
                
                # Fetch as tar stream from source
                bits, stat = src_cont.get_archive(f"{source_dir}/{file_name}")
                tar_data = io.BytesIO()
                for chunk in bits: tar_data.write(chunk)
                tar_data.seek(0)
                
                if target_cont.put_archive(target_dir, tar_data):
                    return f"Moved '{file_name}' from {source_env} to {target_env}."

        except Exception as e:
            logger.exception("Error caught: %s", e)
            return f"Move failed: {str(e)}"

    elif action == "project":
        if not file_name: return "Error: 'file_name' required."
        
        # Correctly identify the source container for projection
        env_to_use = source_env if (source_env and source_env != "chat") else target_env
        
        source_name = get_container_name(env_to_use)
        source_dir = "/lab" if env_to_use == "lab" else "/app"
        
        try:
            container = client.containers.get(source_name)
            
            # Check if it is a directory
            check_dir = container.exec_run(f"test -d {source_dir}/{file_name}")
            is_dir = (check_dir.exit_code == 0)
            
            final_file_name = file_name
            if is_dir:
                # Zip it in the container first
                zip_name = f"{file_name}_export.tar.gz"
                container.exec_run(f"tar -czf /tmp/{zip_name} -C {source_dir}/{file_name} .")
                src_full_path = f"/tmp/{zip_name}"
                final_file_name = zip_name
            else:
                src_full_path = f"{source_dir}/{file_name}"

            # Get archive from container
            bits, stat = container.get_archive(src_full_path)
            tar_data = io.BytesIO()
            for chunk in bits: tar_data.write(chunk)
            tar_data.seek(0)
            
            # Extract from tar to get raw bytes
            with tarfile.open(fileobj=tar_data) as tar:
                member = tar.getmembers()[0]
                file_bytes = tar.extractfile(member).read()

            output_dir = "outputs"
            os.makedirs(output_dir, exist_ok=True)
            unique_name = f"proj_{uuid.uuid4().hex[:6]}_{final_file_name}"
            with open(os.path.join(output_dir, unique_name), "wb") as f:
                f.write(file_bytes)
                
            return f"Projected successfully: [View {final_file_name}](/view/{unique_name})"
            
        except Exception as e:
            logger.exception("Error caught: %s", e)
            return f"Projection failed: {str(e)}"
            
    return "Error: Invalid action."
            
    return "Error: Invalid action. Use 'read', 'move', or 'project'."


available_tools = [
    request_user_interaction,
    web_search,
    send_self_email,
    schedule_task,
    generate_image,
    make_presentation,
    regenerate_presentation_slide,
    analyze_youtube_video,
    manage_files,
    repo_control,
    lab_execute,
    read_tool_output,
    logs_and_preferences,

    report_process_issue,
    obtain_talent,
    compress_memory
]

def __getattr__(name):
    """
    Lazy module attribute resolution to support mock patching in unit tests.
    """
    if name == 'TavilyClient':
        from tavily import TavilyClient
        return TavilyClient
    if name == 'genai':
        from google import genai
        return genai
    if name == 'types':
        from google.genai import types
        return types
    if name == 'requests':
        import requests
        return requests
    raise AttributeError(f"module {__name__} has no attribute {name}")
